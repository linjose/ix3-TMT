#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
pptx_to_md.py
--------------
將 .pptx 轉為 Markdown 風格的 .txt 檔：
- 依頁碼與文字在投影片上的位置順序擷取內容（以 Top/Left 近似閱讀順序）
- 自動從全文擷取關鍵字，作為 tags，輸出在 Markdown 第一部分
- 第二部分起為各投影片（依序）之文字內容
- 可選擇是否包含投影片備註（speaker notes）
- 支援表格輸出為 Markdown 表格或以清單呈現

依賴：python-pptx（pip install python-pptx）
（可選）jieba（若安裝，中文關鍵詞擷取會更準；未安裝則使用 n-gram 後備機制）

使用範例：
python pptx_to_md.py input.pptx -o output.txt --max-tags 20 --include-notes --table-format table
"""

import argparse
import re
import sys
from collections import Counter
from pathlib import Path
from typing import List, Tuple

try:
    from pptx import Presentation
    from pptx.shapes.base import BaseShape
except Exception as e:
    print("請先安裝 python-pptx： pip install python-pptx", file=sys.stderr)
    raise

# ------------------------------
# 基本停用詞（中英）
# ------------------------------

EN_STOPWORDS = set("""
a an and are as at be but by for if in into is it no not of on or s such t that the their then there these they this to was will with you your from we our us
""".split())

# 精簡版中文停用詞（可自行擴充）
ZH_STOPWORDS = set(list("的一是在不了有和就都而及與著或也很到說要與及為於把被這那能個們上中下後前來去及並與等每及及其則於並且"))

# 嘗試載入 jieba（若不可用則退回簡易 n-gram）
try:
    import jieba  # type: ignore
    _HAS_JIEBA = True
except Exception:
    _HAS_JIEBA = False

CJK_RANGE = (
    ("\u4e00", "\u9fff"),   # CJK Unified Ideographs
    ("\u3400", "\u4dbf"),   # CJK Extension A
    ("\uF900", "\uFAFF"),   # CJK Compatibility Ideographs
)

def has_cjk(text: str) -> bool:
    for ch in text:
        for lo, hi in CJK_RANGE:
            if lo <= ch <= hi:
                return True
    return False

def tokenize_en(text: str) -> List[str]:
    words = re.findall(r"[A-Za-z][A-Za-z0-9_'-]*", text.lower())
    return [w for w in words if w not in EN_STOPWORDS and len(w) >= 2]

def tokenize_zh(text: str) -> List[str]:
    # 優先使用 jieba；否則以 2-3 字 n-gram 粗略擷取詞片段
    chars = re.findall(r"[\u3400-\u9fff]", text)
    joined = "".join(chars)
    if not joined:
        return []
    if _HAS_JIEBA:
        segs = [w.strip() for w in jieba.cut(joined) if w.strip()]
        return [w for w in segs if w not in ZH_STOPWORDS and len(w) >= 2]
    else:
        tokens = []
        for n in (2, 3):
            tokens += [joined[i:i+n] for i in range(len(joined)-n+1)]
        tokens = [t for t in tokens if t not in ZH_STOPWORDS]
        return tokens

def extract_keywords(text: str, topk: int = 20) -> List[str]:
    # 依據是否含中日韓文字採用不同 tokenizer；若都沒有，則走英文路徑
    tokens = []
    if has_cjk(text):
        tokens += tokenize_zh(text)
    # 也同時擷取英文（投影片常見中英混排）
    tokens += tokenize_en(text)
    if not tokens:
        return []
    ctr = Counter(tokens)
    # 依頻率排序，長度較長者略有優先（同頻時）
    items = sorted(ctr.items(), key=lambda kv: (kv[1], len(kv[0])), reverse=True)
    return [w for w, _ in items[:topk]]

# ------------------------------
# 形狀與段落文字擷取
# ------------------------------

def _shape_bounds(shape: BaseShape) -> Tuple[int, int]:
    # 使用 top、left 作為排序鍵（EMU 單位）；個別 shape 可能缺屬性，故保護性處理
    top = getattr(shape, "top", 0)
    left = getattr(shape, "left", 0)
    try:
        return int(top), int(left)
    except Exception:
        return (0, 0)

def _escape_md(text: str) -> str:
    # 簡易 Markdown 字元跳脫
    return (text
            .replace("\\", "\\\\")
            .replace("|", "\\|")
            .replace("_", "\\_")
            .replace("*", "\\*")
            .replace("`", "\\`")
            .replace("#", "\\#"))

def _paragraphs_to_md(paragraphs) -> List[str]:
    lines = []
    for p in paragraphs:
        txt = "".join(run.text for run in p.runs) if getattr(p, "runs", None) else p.text
        if not txt or not txt.strip():
            continue
        level = getattr(p, "level", 0) or 0
        indent = "  " * level
        # 以無序清單表示（無法可靠取用編號資訊，因此一律以 - ）
        lines.append(f"{indent}- {_escape_md(txt.strip())}")
    return lines

def _table_to_md(table, table_format: str = "table") -> List[str]:
    if table_format not in ("table", "list"):
        table_format = "list"
    rows = []
    for r in table.rows:
        cells = []
        for c in r.cells:
            # 單一儲存格內可能有多段，合併為單一行
            texts = []
            if hasattr(c, "text_frame") and c.text_frame:
                for p in c.text_frame.paragraphs:
                    if getattr(p, "text", "").strip():
                        texts.append(p.text.strip())
            else:
                texts.append(c.text.strip())
            cell_text = " ".join(texts).strip()
            cells.append(_escape_md(cell_text))
        rows.append(cells)

    if not rows:
        return []

    if table_format == "list":
        out = []
        for ridx, cells in enumerate(rows, start=1):
            out.append(f"- 表格第 {ridx} 行： " + " | ".join(cells))
        return out

    # Markdown 表格
    # 若只有一行，當作一般列輸出
    if len(rows) == 1:
        return [f"- 表格： " + " | ".join(rows[0])]

    header = rows[0]
    body = rows[1:]
    lines = []
    lines.append("| " + " | ".join(header) + " |")
    lines.append("| " + " | ".join(["---"] * len(header)) + " |")
    for r in body:
        lines.append("| " + " | ".join(r) + " |")
    return lines

def shape_text_to_md(shape: BaseShape, table_format: str = "list") -> List[str]:
    # 文本框、標題、placeholder
    if hasattr(shape, "has_text_frame") and shape.has_text_frame:
        return _paragraphs_to_md(shape.text_frame.paragraphs)

    # 表格
    # 19 對應 MSO_SHAPE_TYPE.TABLE（避免額外依賴列舉）
    if getattr(shape, "shape_type", None) == 19 and hasattr(shape, "table"):
        return _table_to_md(shape.table, table_format=table_format)

    # 其他類型（圖片、圖表、SmartArt 等）以占位文字表示
    if getattr(shape, "shape_type", None) is not None:
        return [f"- [物件：{getattr(shape, 'shape_type', 'Unknown')}]" ]

    return []

def extract_slides_text(prs: Presentation, include_notes: bool = False, table_format: str = "list") -> Tuple[List[List[str]], str]:
    """
    回傳：
    - slides_lines: 每張投影片的 Markdown 行（List[str]）
    - all_text: 全文（便於抽取關鍵字）
    """
    slides_lines: List[List[str]] = []
    corpus_parts: List[str] = []

    for s_idx, slide in enumerate(prs.slides, start=1):
        # 按「上到下、左到右」排序 shape，以近似視覺閱讀順序
        shapes_sorted = sorted(slide.shapes, key=_shape_bounds)

        lines: List[str] = []
        # 投影片標題（若有）
        title = getattr(slide.shapes.title, "text", None)
        if title and title.strip():
            lines.append(f"### 第 {s_idx} 頁：{_escape_md(title.strip())}")
            corpus_parts.append(title.strip())
        else:
            lines.append(f"### 第 {s_idx} 頁")

        for sh in shapes_sorted:
            try:
                md_lines = shape_text_to_md(sh, table_format=table_format)
                for ln in md_lines:
                    if ln.strip():
                        lines.append(ln)
                        corpus_parts.append(re.sub(r"^- ", "", ln).strip())
            except Exception:
                # 遇到不支援形狀，略過
                continue

        # 備註
        if include_notes and getattr(slide, "has_notes_slide", False):
            notes_slide = slide.notes_slide
            if notes_slide and getattr(notes_slide, "notes_text_frame", None):
                note_lines = _paragraphs_to_md(notes_slide.notes_text_frame.paragraphs)
                if note_lines:
                    lines.append("#### 備註")
                    lines.extend(note_lines)
                    for nl in note_lines:
                        corpus_parts.append(re.sub(r"^- ", "", nl).strip())

        # 如果除標題外無內容，加入空白提示
        if len(lines) == 1:
            lines.append("- （本頁無可擷取的文字內容）")

        slides_lines.append(lines)

    all_text = "\n".join(corpus_parts)
    return slides_lines, all_text

def build_markdown(tags: List[str], slides_lines: List[List[str]]) -> str:
    md = []
    md.append("# Tags")
    if tags:
        for t in tags:
            md.append(f"- {t}")
    else:
        md.append("- （無）")

    md.append("\n---\n")
    md.append("# 內容")

    for lines in slides_lines:
        md.extend(lines)
        md.append("")  # 分隔空行

    return "\n".join(md).strip() + "\n"

def convert(pptx_path: Path, out_path: Path, max_tags: int, include_notes: bool, table_format: str):
    prs = Presentation(str(pptx_path))
    slides_lines, all_text = extract_slides_text(prs, include_notes=include_notes, table_format=table_format)
    tags = extract_keywords(all_text, topk=max_tags)
    md = build_markdown(tags, slides_lines)
    out_path.write_text(md, encoding="utf-8")
    print(f"✅ 已輸出：{out_path}")

def main():
    p = argparse.ArgumentParser(description="將 PPTX 轉為 Markdown 風格的 TXT（含標籤 Tags）。")
    p.add_argument("pptx", type=str, help="輸入的 .pptx 檔案路徑")
    p.add_argument("-o", "--output", type=str, default=None, help="輸出的 .txt（markdown）檔案路徑（預設為同名 .txt）")
    p.add_argument("--max-tags", type=int, default=20, help="輸出關鍵字（tags）數量上限（預設 20）")
    p.add_argument("--include-notes", action="store_true", help="是否將投影片備註（speaker notes）一併輸出")
    p.add_argument("--table-format", choices=["list", "table"], default="list", help="表格輸出格式：list 或 table（預設 list）")

    args = p.parse_args()

    pptx_path = Path(args.pptx)
    if not pptx_path.exists():
        print(f"找不到檔案：{pptx_path}", file=sys.stderr)
        sys.exit(1)

    out_path = Path(args.output) if args.output else pptx_path.with_suffix(".txt")
    try:
        convert(pptx_path, out_path, max_tags=args.max_tags, include_notes=args.include_notes, table_format=args.table_format)
    except Exception as e:
        print(f"轉換失敗：{e}", file=sys.stderr)
        sys.exit(2)

if __name__ == "__main__":
    main()
